"""
Script para generar una presentación PowerPoint con el contenido y estilo de iTTel
Basado en el sitio web index2.html
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os

def crear_presentacion_ittel():
    """Crea una presentación con el look and feel de la web de iTTel"""
    
    # Crear presentación
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Colores basados en el sitio web
    COLOR_PRIMARY = RGBColor(56, 189, 248)  # #38bdf8 - azul principal
    COLOR_TEXT_DARK = RGBColor(17, 24, 39)  # #111827 - texto oscuro
    COLOR_TEXT_MEDIUM = RGBColor(107, 114, 128)  # #6b7280 - texto medio
    COLOR_WHITE = RGBColor(255, 255, 255)
    COLOR_BG_LIGHT = RGBColor(248, 250, 252)  # #f8fafc - fondo claro
    
    # --- SLIDE 1: Portada ---
    slide_portada = prs.slides.add_slide(prs.slide_layouts[6])  # Layout en blanco
    
    # Fondo
    background = slide_portada.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_PRIMARY
    
    # Logo (si existe, agregarlo manualmente o usar texto)
    logo_path = os.path.join("assets", "images", "logoIttel1.jpg")
    if os.path.exists(logo_path):
        left = Inches(0.5)
        top = Inches(0.5)
        logo = slide_portada.shapes.add_picture(logo_path, left, top, width=Inches(2))
    
    # Título principal
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(1.5)
    
    title_box = slide_portada.shapes.add_textbox(left, top, width, height)
    text_frame = title_box.text_frame
    text_frame.word_wrap = True
    
    p = text_frame.paragraphs[0]
    p.text = "Grupo iTTel"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtítulo
    left = Inches(1)
    top = Inches(4)
    width = Inches(8)
    height = Inches(1)
    
    subtitle_box = slide_portada.shapes.add_textbox(left, top, width, height)
    text_frame = subtitle_box.text_frame
    
    p = text_frame.paragraphs[0]
    p.text = "Information Technology & Telecommunications"
    p.font.size = Pt(28)
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # --- SLIDE 2: Nosotros - Intro ---
    slide_nosotros = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fondo blanco
    background = slide_nosotros.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_WHITE
    
    # Subtítulo superior
    left = Inches(1)
    top = Inches(0.8)
    width = Inches(8)
    height = Inches(0.5)
    
    subtitle_box = slide_nosotros.shapes.add_textbox(left, top, width, height)
    p = subtitle_box.text_frame.paragraphs[0]
    p.text = "GRUPO ITTEL"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.alignment = PP_ALIGN.LEFT
    
    # Título
    left = Inches(1)
    top = Inches(1.3)
    width = Inches(8)
    height = Inches(0.8)
    
    title_box = slide_nosotros.shapes.add_textbox(left, top, width, height)
    p = title_box.text_frame.paragraphs[0]
    p.text = "Conectando el futuro, hoy."
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_DARK
    p.alignment = PP_ALIGN.LEFT
    
    # Descripción
    left = Inches(1)
    top = Inches(2.3)
    width = Inches(8)
    height = Inches(1.2)
    
    desc_box = slide_nosotros.shapes.add_textbox(left, top, width, height)
    text_frame = desc_box.text_frame
    text_frame.word_wrap = True
    
    p = text_frame.paragraphs[0]
    p.text = "Somos una Empresa con vasta experiencia en el mercado de las telecomunicaciones, que reconoce el impacto y el Valor Agregado que la innovación tecnológica genera en las organizaciones y comunidades."
    p.font.size = Pt(18)
    p.font.color.rgb = COLOR_TEXT_MEDIUM
    p.line_spacing = 1.5
    
    # --- SLIDE 3: Misión, Visión y Valores ---
    slide_mvv = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fondo
    background = slide_mvv.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG_LIGHT
    
    # Título
    left = Inches(1)
    top = Inches(0.5)
    width = Inches(8)
    height = Inches(0.7)
    
    title_box = slide_mvv.shapes.add_textbox(left, top, width, height)
    p = title_box.text_frame.paragraphs[0]
    p.text = "Nuestra Esencia"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_DARK
    p.alignment = PP_ALIGN.CENTER
    
    # Crear tres cajas para Misión, Visión, Valores
    cards = [
        {
            "icon": "🚀",
            "title": "Nuestra Misión",
            "desc": "Brindar soluciones integrales en tecnología y telecomunicaciones, adaptadas a las necesidades de cada cliente, garantizando calidad, eficiencia y un servicio de postventa excepcional."
        },
        {
            "icon": "👁",
            "title": "Nuestra Visión",
            "desc": "Ser la empresa de tecnología y telecomunicaciones líder en la región, reconocida por nuestra innovación, confiabilidad y el compromiso con el éxito de nuestros clientes."
        },
        {
            "icon": "💎",
            "title": "Nuestros Valores",
            "desc": "Compromiso, Calidad, Innovación, Confianza y Orientación al cliente. Estos pilares guían cada uno de nuestros proyectos y relaciones comerciales."
        }
    ]
    
    card_top = 1.5
    card_height = 1.6
    
    for i, card in enumerate(cards):
        left = Inches(1)
        top = Inches(card_top + (i * 1.85))
        width = Inches(8)
        height = Inches(card_height)
        
        # Crear rectángulo para la tarjeta
        shape = slide_mvv.shapes.add_shape(
            1,  # Rectangle
            left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_WHITE
        shape.line.color.rgb = RGBColor(229, 231, 235)
        shape.shadow.inherit = False
        
        # Icono y título
        text_frame = shape.text_frame
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0.3)
        text_frame.margin_right = Inches(0.3)
        text_frame.margin_top = Inches(0.2)
        
        p = text_frame.paragraphs[0]
        p.text = f"{card['icon']} {card['title']}"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_DARK
        
        # Descripción
        p2 = text_frame.add_paragraph()
        p2.text = card['desc']
        p2.font.size = Pt(14)
        p2.font.color.rgb = COLOR_TEXT_MEDIUM
        p2.space_before = Pt(8)
        p2.line_spacing = 1.3
    
    # --- SLIDE 4: IT & Telco - Pioneros ---
    slide_ittelco = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fondo
    background = slide_ittelco.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_WHITE
    
    # Título sección
    left = Inches(1)
    top = Inches(0.5)
    width = Inches(8)
    height = Inches(0.7)
    
    title_box = slide_ittelco.shapes.add_textbox(left, top, width, height)
    p = title_box.text_frame.paragraphs[0]
    p.text = "IT & Telco"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_DARK
    p.alignment = PP_ALIGN.CENTER
    
    # Subtítulo
    left = Inches(1)
    top = Inches(1.2)
    width = Inches(8)
    height = Inches(0.5)
    
    subtitle_box = slide_ittelco.shapes.add_textbox(left, top, width, height)
    p = subtitle_box.text_frame.paragraphs[0]
    p.text = "Nuestra experiencia y desarrollos nos permiten ofrecer un servicio diferencial en el mercado."
    p.font.size = Pt(16)
    p.font.color.rgb = COLOR_TEXT_MEDIUM
    p.alignment = PP_ALIGN.CENTER
    
    # Tarjeta destacada
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(2.5)
    
    shape = slide_ittelco.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(240, 249, 255)
    shape.line.color.rgb = COLOR_PRIMARY
    shape.line.width = Pt(2)
    
    text_frame = shape.text_frame
    text_frame.word_wrap = True
    text_frame.margin_left = Inches(0.4)
    text_frame.margin_right = Inches(0.4)
    text_frame.margin_top = Inches(0.3)
    
    p = text_frame.paragraphs[0]
    p.text = "🌐 Pioneros en Redes Neutrales"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_DARK
    
    p2 = text_frame.add_paragraph()
    p2.text = "Fuimos los primeros en Argentina en implementar lo que hoy se conoce como 'Redes Neutrales' o 'Compartición de Infraestructura Telco'. Llevamos más de 10 años perfeccionando nuestros desarrollos en las más variadas jurisdicciones."
    p2.font.size = Pt(16)
    p2.font.color.rgb = COLOR_TEXT_MEDIUM
    p2.space_before = Pt(12)
    p2.line_spacing = 1.4
    
    # Características en la parte inferior
    features = [
        "✓ Convergencia Total",
        "✓ Sinergia Operativa",
        "✓ SLA Premium",
        "✓ Implementación Sustentable"
    ]
    
    feature_left = 1
    feature_top = 4.8
    feature_width = 1.8
    
    for i, feature in enumerate(features):
        left = Inches(feature_left + (i * 2.05))
        top = Inches(feature_top)
        width = Inches(feature_width)
        height = Inches(0.5)
        
        box = slide_ittelco.shapes.add_textbox(left, top, width, height)
        p = box.text_frame.paragraphs[0]
        p.text = feature
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        p.alignment = PP_ALIGN.CENTER
    
    # --- SLIDE 5: Estadísticas ---
    slide_stats = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fondo con gradiente (simulado con color sólido)
    background = slide_stats.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_PRIMARY
    
    # Título
    left = Inches(1)
    top = Inches(1)
    width = Inches(8)
    height = Inches(0.7)
    
    title_box = slide_stats.shapes.add_textbox(left, top, width, height)
    p = title_box.text_frame.paragraphs[0]
    p.text = "Números que nos respaldan"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Estadísticas
    stats = [
        {"number": "10+", "label": "Años de\nexperiencia"},
        {"number": "340+", "label": "Proyectos\nde éxito"},
        {"number": "99.99%", "label": "de SLA\ngarantizado"}
    ]
    
    stat_width = 2.3
    stat_left_start = 1.2
    stat_top = 3
    
    for i, stat in enumerate(stats):
        left = Inches(stat_left_start + (i * 2.7))
        top = Inches(stat_top)
        width = Inches(stat_width)
        height = Inches(2)
        
        box = slide_stats.shapes.add_textbox(left, top, width, height)
        text_frame = box.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # Número
        p = text_frame.paragraphs[0]
        p.text = stat['number']
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.alignment = PP_ALIGN.CENTER
        
        # Label
        p2 = text_frame.add_paragraph()
        p2.text = stat['label']
        p2.font.size = Pt(16)
        p2.font.color.rgb = COLOR_WHITE
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(12)
    
    # --- SLIDES 6-12: Proyectos Destacados ---
    proyectos = [
        {
            "title": "AUSA",
            "subtitle": "Acuerdo iTTel – AUSA",
            "brief": "Gestión integral de infraestructura TELCO de alto desempeño en CABA.",
            "details": "• Gestión de más de 1.000 km de hilos arrendados\n• 100+ estructuras de antenas distribuidas\n• SLA > 99,9% con atención 24×7×365\n• Diseño, implementación y operación bajo estándares de alta disponibilidad",
            "image": "AUSA4.jpg"
        },
        {
            "title": "AUBASA",
            "subtitle": "Acuerdo iTTel – AUBASA",
            "brief": "Red de fibra óptica redundante y backbone estratégico CABA - La Plata.",
            "details": "• Red de fibra óptica con topología redundante\n• Backbone estratégico entre CABA y La Plata\n• Fibra Oscura entre CABASE y NAP La Plata\n• SLA contractual > 99,99%\n• Cuadrilla in situ con tiempo de respuesta inmediato",
            "image": "AUBASA1.jpg"
        },
        {
            "title": "USITTEL",
            "subtitle": "Red FTTH en Tandil",
            "brief": "Despliegue de red FTTH para Internet y TV en alianza con La Usina de Tandil.",
            "details": "• Arquitectura GPON para Internet y TV digital\n• Cobertura urbana y periurbana\n• SLA premium orientado al usuario final\n• Operación centralizada con monitoreo proactivo\n• Web: usittel.com.ar",
            "image": "proyectos/usittel_internet.png"
        },
        {
            "title": "LARANET",
            "subtitle": "Red FTTH en La Matanza",
            "brief": "Diseño e implementación de red óptica pasiva con alta capacidad de crecimiento.",
            "details": "• Red óptica pasiva FTTH en La Matanza\n• Arquitectura escalable con backbone redundante\n• Operación 24/7 desde NOC\n• Implementación por etapas\n• Web: laranet.com.ar",
            "image": "proyectos/laranett.png"
        },
        {
            "title": "ADIFSE",
            "subtitle": "Proyecto ADIFSE",
            "brief": "Relevamiento de infraestructura en trazados ferroviarios.",
            "details": "• Primera etapa de relevamiento de infraestructura\n• Líneas: Mitre, San Martín, Sarmiento, Belgrano Sur, Gral. Roca y Tren de la Costa\n• Base para planificar desarrollo sustentable de redes\n• Optimización de operación e infraestructura TELCO",
            "image": "proyectos/adifse_portada.png"
        },
        {
            "title": "VITTAL",
            "subtitle": "Videovigilancia",
            "brief": "Sistema complejo de videovigilancia IP con grabación y acceso centralizado.",
            "details": "• Plataforma de Video Vigilancia offline y online\n• Accesible desde cualquier dispositivo con Internet\n• Visualización y grabación centralizada\n• Control remoto de cámaras\n• Soporte 4K, visión nocturna, detección por movimiento",
            "image": "proyectos/viTTal.png"
        },
        {
            "title": "CARREFOUR",
            "subtitle": "Infraestructura IT",
            "brief": "Infraestructura tecnológica integral con conectividad y monitoreo 24/7.",
            "details": "• Solución integral para auto-servicio de comedor\n• Validación por huella digital o tarjeta de proximidad\n• Múltiples menús con control automático\n• Registro de transacciones con reportes Excel\n• Auditoría centralizada y escalable",
            "image": "proyectos/carrefour.png"
        }
    ]
    
    for proyecto in proyectos:
        slide_proyecto = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Fondo blanco
        background = slide_proyecto.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_WHITE
        
        # Título del proyecto (grande)
        left = Inches(0.5)
        top = Inches(0.4)
        width = Inches(4)
        height = Inches(0.8)
        
        title_box = slide_proyecto.shapes.add_textbox(left, top, width, height)
        p = title_box.text_frame.paragraphs[0]
        p.text = proyecto['title']
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        
        # Subtítulo
        left = Inches(0.5)
        top = Inches(1.2)
        width = Inches(4)
        height = Inches(0.5)
        
        subtitle_box = slide_proyecto.shapes.add_textbox(left, top, width, height)
        p = subtitle_box.text_frame.paragraphs[0]
        p.text = proyecto['subtitle']
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_DARK
        
        # Brief
        left = Inches(0.5)
        top = Inches(1.8)
        width = Inches(4)
        height = Inches(0.7)
        
        brief_box = slide_proyecto.shapes.add_textbox(left, top, width, height)
        text_frame = brief_box.text_frame
        text_frame.word_wrap = True
        p = text_frame.paragraphs[0]
        p.text = proyecto['brief']
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_TEXT_MEDIUM
        p.line_spacing = 1.3
        
        # Detalles
        left = Inches(0.5)
        top = Inches(2.7)
        width = Inches(4)
        height = Inches(4)
        
        details_box = slide_proyecto.shapes.add_textbox(left, top, width, height)
        text_frame = details_box.text_frame
        text_frame.word_wrap = True
        p = text_frame.paragraphs[0]
        p.text = proyecto['details']
        p.font.size = Pt(13)
        p.font.color.rgb = COLOR_TEXT_DARK
        p.line_spacing = 1.5
        
        # Agregar imagen si existe
        image_path = os.path.join("assets", "images", proyecto['image'])
        if os.path.exists(image_path):
            left = Inches(5)
            top = Inches(1.5)
            try:
                pic = slide_proyecto.shapes.add_picture(
                    image_path, left, top, height=Inches(4.5)
                )
            except:
                # Si hay error con la imagen, continuar sin ella
                pass
    
    # --- SLIDE FINAL: Contacto ---
    slide_contacto = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fondo
    background = slide_contacto.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG_LIGHT
    
    # Título
    left = Inches(1)
    top = Inches(0.8)
    width = Inches(8)
    height = Inches(0.8)
    
    title_box = slide_contacto.shapes.add_textbox(left, top, width, height)
    p = title_box.text_frame.paragraphs[0]
    p.text = "Contactanos"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_DARK
    p.alignment = PP_ALIGN.CENTER
    
    # Información de contacto
    contacto_info = [
        {
            "icon": "📞",
            "title": "Teléfonos",
            "info": "011-5272-9000\n0810-345-ITTEL (4883)"
        },
        {
            "icon": "📧",
            "title": "Emails",
            "info": "administracion@it-tel.com.ar\nsoporte@it-tel.com.ar"
        },
        {
            "icon": "💼",
            "title": "LinkedIn",
            "info": "/company/grupo-ittel"
        }
    ]
    
    contact_top = 2.2
    contact_height = 1.5
    
    for i, contact in enumerate(contacto_info):
        left = Inches(1)
        top = Inches(contact_top + (i * 1.7))
        width = Inches(8)
        height = Inches(contact_height)
        
        # Crear rectángulo
        shape = slide_contacto.shapes.add_shape(1, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_WHITE
        shape.line.color.rgb = RGBColor(229, 231, 235)
        
        text_frame = shape.text_frame
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0.3)
        text_frame.margin_top = Inches(0.2)
        
        # Título
        p = text_frame.paragraphs[0]
        p.text = f"{contact['icon']} {contact['title']}"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_DARK
        
        # Info
        p2 = text_frame.add_paragraph()
        p2.text = contact['info']
        p2.font.size = Pt(16)
        p2.font.color.rgb = COLOR_TEXT_MEDIUM
        p2.space_before = Pt(8)
    
    # Guardar presentación
    output_path = "Presentacion_iTTel.pptx"
    prs.save(output_path)
    print(f"✅ Presentación creada exitosamente: {output_path}")
    print(f"📊 Total de slides: {len(prs.slides)}")
    return output_path

if __name__ == "__main__":
    print("🚀 Generando presentación PowerPoint para iTTel...")
    print("=" * 60)
    
    try:
        archivo = crear_presentacion_ittel()
        print("=" * 60)
        print(f"✨ ¡Listo! Tu presentación está en: {archivo}")
        print("\n📝 La presentación incluye:")
        print("   • Portada con branding iTTel")
        print("   • Sección Nosotros (Misión, Visión, Valores)")
        print("   • IT & Telco con características destacadas")
        print("   • Estadísticas corporativas")
        print("   • 7 Proyectos destacados con detalles")
        print("   • Slide de contacto")
        print("\n💡 Recuerda verificar que las imágenes estén en la carpeta 'assets/images'")
    except Exception as e:
        print(f"❌ Error al generar la presentación: {e}")
        print("\n💡 Asegúrate de tener instalada la librería python-pptx:")
        print("   pip install python-pptx")
