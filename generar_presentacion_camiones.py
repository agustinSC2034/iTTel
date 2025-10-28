"""
Genera una presentación PowerPoint profesional para la propuesta de
Videovigilancia Móvil de la flota de recolección de residuos del GCBA,
usando el look & feel del sitio (paleta y estilo) y contenido técnico
desde `presentacion_camiones/contexto2.md`.

Requiere: python-pptx, Pillow
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


# Paleta y estilo tomados de assets/css/main.css (:root)
COLOR_PRIMARY = RGBColor(34, 125, 179)   # #227db3
COLOR_ACCENT = RGBColor(56, 189, 248)    # #38bdf8
COLOR_SECONDARY = RGBColor(0, 212, 255)  # #00d4ff
COLOR_TEXT_DARK = RGBColor(26, 32, 44)   # #1a202c
COLOR_TEXT_MEDIUM = RGBColor(74, 85, 104)  # #4a5568
COLOR_TEXT_LIGHT = RGBColor(100, 116, 139)  # #64748b
COLOR_BG_LIGHT = RGBColor(248, 250, 252)  # #f8fafc
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_BORDER = RGBColor(229, 231, 235)


def add_title(slide, title, subtitle=None, title_color=COLOR_WHITE, center=True):
    """Utilidad para agregar título (y opcional subtítulo)."""
    # Título
    tb = slide.shapes.add_textbox(Inches(0.75), Inches(2.2), Inches(8.5), Inches(1.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(46)
    p.font.bold = True
    p.font.color.rgb = title_color
    p.alignment = PP_ALIGN.CENTER if center else PP_ALIGN.LEFT

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.color.rgb = COLOR_WHITE if title_color == COLOR_WHITE else COLOR_TEXT_MEDIUM
        p2.space_before = Pt(10)
        p2.alignment = PP_ALIGN.CENTER if center else PP_ALIGN.LEFT


def add_section_header(slide, text):
    """Encabezado de sección con estilo web."""
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(8.8), Inches(0.8))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_DARK
    p.alignment = PP_ALIGN.LEFT


def add_bullets(slide, items, left=0.9, top=1.5, width=8.2, height=4.8, font_size=18, color=COLOR_TEXT_DARK):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.line_spacing = 1.35


def add_card(slide, title, description, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_WHITE
    shape.line.color.rgb = COLOR_BORDER
    shape.shadow.inherit = False

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    tf.margin_right = Inches(0.3)
    tf.margin_top = Inches(0.2)

    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_DARK

    p2 = tf.add_paragraph()
    p2.text = description
    p2.font.size = Pt(14)
    p2.font.color.rgb = COLOR_TEXT_MEDIUM
    p2.space_before = Pt(6)
    p2.line_spacing = 1.35


def safe_add_picture(slide, img_path, left, top, width=None, height=None):
    if os.path.exists(img_path):
        try:
            if width:
                slide.shapes.add_picture(img_path, Inches(left), Inches(top), width=Inches(width))
            elif height:
                slide.shapes.add_picture(img_path, Inches(left), Inches(top), height=Inches(height))
            else:
                slide.shapes.add_picture(img_path, Inches(left), Inches(top))
        except Exception:
            pass


def generar_presentacion_camiones():
    # Intentar reutilizar el estilo de una presentación corporativa previa como plantilla
    template_path = os.path.join(os.getcwd(), "Presentacion_Grupo-iTTel.pptx")
    try:
        prs = Presentation(template_path) if os.path.exists(template_path) else Presentation()
    except Exception:
        prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    assets = os.path.join("assets", "images")
    camiones_dir = os.path.join("presentacion_camiones")
    logo = os.path.join(assets, "logoIttel1.jpg")

    # Slide 1: Portada
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s1.background.fill
    bg.solid()
    bg.fore_color.rgb = COLOR_PRIMARY
    # Logo
    safe_add_picture(s1, logo, left=0.5, top=0.5, width=2.0)
    add_title(
        s1,
        title="Propuesta Comercial",
        subtitle="Videovigilancia Móvil para Flota de Residuos – GCBA",
        title_color=COLOR_WHITE,
        center=True,
    )

    # Slide 2: Resumen Ejecutivo
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(s2, "Resumen ejecutivo")
    bullets = [
        "Objetivo: monitoreo, auditoría y seguridad 24/7 de la flota.",
        "Ecosistema MDVR + cámaras + comunicación 4G/GPS + VMS central (Hik-Central).",
        "Auditoría acelerada por eventos (sensor en brazo hidráulico) y G-sensor.",
        "Equipamiento de grado vehicular (vibración/agua/polvo) y soporte 24×7×365.",
    ]
    add_bullets(s2, bullets)
    safe_add_picture(s2, os.path.join(assets, "proyectos", "solucion_VITTAL.png"), left=6.6, top=1.1, height=3.6)

    # Slide 3: El desafío
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(s3, "El desafío")
    bullets = [
        "Auditar el proceso de levantamiento de contenedores para reclamos ciudadanos.",
        "Generar evidencia irrefutable ante incidentes viales y operativos.",
        "Gestión de ubicación, eventos y estado de la flota en tiempo real.",
        "Entorno hostil: vibración constante, agua, polvo y uso rudo 24/7.",
    ]
    add_bullets(s3, bullets)

    # Slide 4: La solución (visión general)
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(s4, "La solución propuesta")
    add_card(s4, "Hardware embarcado (MDVR)",
             "Grabación multi‑canal, módulos 4G/GPS, G‑sensor y supercapacitores.", 0.6, 1.4, 2.9, 1.8)
    add_card(s4, "Red de comunicación",
             "Transmisión en vivo y telemetría por 4G/LTE; GPS para geoposición.", 3.6, 1.4, 2.9, 1.8)
    add_card(s4, "Software de gestión",
             "Hik‑Central Professional: mapa GIS, live view, búsqueda por eventos, protección de evidencia.", 6.6, 1.4, 2.9, 1.8)

    # Slide 5: Equipamiento – MDVR
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(s5, "Equipamiento – MDVR")
    add_bullets(s5, [
        "Hikvision DS‑MP5604H‑S/GW (híbrido analógico/IP).",
        "Módulos 4G/LTE y GPS integrados (/GW).",
        "G‑Sensor integrado (colisiones/frenadas/vuelcos).",
        "Supercapacitores: cierre seguro ante pérdida de batería.",
        "Bahía 2.5\" con llave para SSD industrial.",
    ], top=1.4)
    # Imagen del MDVR (fondo transparente)
    safe_add_picture(
        s5,
        os.path.join(camiones_dir, "MVR DS-MP5604-SD.png"),
        left=6.4,
        top=1.1,
        height=3.9,
    )

    # Slide 6: Almacenamiento SSD
    s6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(s6, "Almacenamiento – SSD industrial")
    add_bullets(s6, [
        "Resistente a vibración y golpes: clave en camiones.",
        "Evita fallas típicas de HDD y tarjetas SD.",
        "Mayor confiabilidad de evidencia y mayor durabilidad.",
    ], top=1.4)

    # Slide 7: Cámaras y conectores
    s7 = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(s7, "Cámaras y conectores")
    add_bullets(s7, [
        "Cámaras Hikvision serie DS‑2CS… (ej. DS‑2CS54D7T‑M).",
        "Antivandálicas (IK10) y herméticas (IP67).",
        "Conectores de aviación M12 roscados para vibración: conexión permanente.",
    ], top=1.4)
    # Imagen de cámara (fondo transparente)
    safe_add_picture(
        s7,
        os.path.join(camiones_dir, "camara_domo.png"),
        left=6.5,
        top=1.0,
        height=4.0,
    )

    # Slide 8: Integración eléctrica e ignición
    s8 = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(s8, "Integración eléctrica e ignición")
    add_bullets(s8, [
        "Instalación oculta en cabina (segura).",
        "Cableado B+ (batería), GND (tierra) e IGN (contacto).",
        "Operación 100% automática: graba con contacto y apaga con retardo programado.",
    ], top=1.4)

    # Slide 9: Sensor de brazo y auditoría por eventos
    s9 = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(s9, "Auditoría: sensor en brazo hidráulico")
    add_bullets(s9, [
        "Sensor externo cableado a Entrada de Alarma del MDVR.",
        "Cada activación genera una marca de evento (TAG) en la grabación.",
        "Búsqueda por evento en Hik‑Central: auditoría en segundos.",
    ], top=1.4)

    # Slide 10: Plataforma de monitoreo
    s10 = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(s10, "Plataforma de monitoreo – Hik‑Central")
    add_bullets(s10, [
        "Mapa GIS con la flota (GPS en tiempo real).",
        "Live view por 4G/LTE y reproducción centralizada.",
        "Protección de evidencia y roles/permiso por perfil.",
        "Clientes web y apps móviles (Android/iOS).",
    ], top=1.4)
    safe_add_picture(s10, os.path.join(assets, "proyectos", "viTTal.png"), left=6.6, top=1.1, height=3.6)

    # Slide 11: Casos de uso
    s11 = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(s11, "Casos de uso clave")
    add_card(s11, "Monitoreo en vivo",
             "Operadores visualizan cámaras y ubicación de camiones en tiempo real.", 0.6, 1.4, 2.9, 1.6)
    add_card(s11, "Auditoría por eventos",
             "Filtrado por TAG de brazo hidráulico; clips de 30 s por evento.", 3.6, 1.4, 2.9, 1.6)
    add_card(s11, "Gestión de incidentes",
             "G‑sensor marca colisiones/frenadas; evidencia protegida automáticamente.", 6.6, 1.4, 2.9, 1.6)

    # Slide 12: Arquitectura de comunicación
    s12 = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(s12, "Arquitectura de comunicación")
    add_bullets(s12, [
        "Camión: Cámaras → MDVR → 4G/LTE + GPS.",
        "Red móvil → Centro de Datos → Hik‑Central.",
        "Operación 24×7×365, monitoreo y reportes.",
    ], top=1.4)

    # Slide 13: Plan de implementación
    s13 = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(s13, "Plan de implementación")
    add_bullets(s13, [
        "Relevamiento técnico y diseño (1–2 semanas).",
        "Piloto en X camiones (2–4 semanas).",
        "Despliegue por lotes (escala progresiva).",
        "Capacitación operadores y SOPs.",
        "Aseguramiento de calidad (UAT) y postventa.",
    ], top=1.4)

    # Slide 14: Beneficios para GCBA
    s14 = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(s14, "Beneficios para GCBA")
    add_bullets(s14, [
        "Mejora de SLA y trazabilidad de servicio.",
        "Reducción de reclamos y tiempos de auditoría.",
        "Seguridad operativa y vial incrementada.",
        "Evidencia forense confiable y protegida.",
        "Datos para analítica y mejora continua.",
    ], top=1.4)

    # Slide 15: SLA y soporte
    s15 = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(s15, "SLA y soporte")
    add_bullets(s15, [
        "Atención 24×7×365 desde NOC.",
        "KPIs operativos y reportes periódicos.",
        "Mantenimiento preventivo y correctivo.",
        "Stock de partes y reemplazo prioritario.",
    ], top=1.4)

    # Slide 16: Experiencia relevante
    s16 = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(s16, "Experiencia relevante")
    add_bullets(s16, [
        "AUSA: Gestión de infraestructura TELCO en CABA (SLA > 99,9%).",
        "AUBASA: Backbone redundante CABA – La Plata (SLA > 99,99%).",
        "VITTAL: Plataforma de videovigilancia IP centralizada.",
    ], top=1.4)
    safe_add_picture(s16, os.path.join(assets, "AUSA4.jpg"), left=6.6, top=1.1, height=1.9)
    safe_add_picture(s16, os.path.join(assets, "proyectos", "ruta_AUBASA.jpg"), left=6.6, top=3.2, height=1.9)

    # Slide 17: Próximos pasos
    s17 = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(s17, "Próximos pasos")
    add_bullets(s17, [
        "Validar alcance y requerimientos (workshop técnico).",
        "Definir cantidad de camiones y canales.",
        "Cronograma piloto y HLD/LLD.",
        "Aprobación y kick‑off de proyecto.",
    ], top=1.4)

    # Slide 18: Contacto
    s18 = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s18.background.fill
    bg.solid()
    bg.fore_color.rgb = COLOR_BG_LIGHT
    add_title(s18, "Contacto", subtitle=None, title_color=COLOR_TEXT_DARK, center=True)

    info = [
        ("Teléfonos", "011-5272-9000\n0810-345-ITTEL (4883)"),
        ("Emails", "administracion@it-tel.com.ar\nsoporte@it-tel.com.ar"),
        ("LinkedIn", "/grupoittel"),
    ]
    top = 3.0
    for title, desc in info:
        add_card(s18, title, desc, left=0.9, top=top, width=8.2, height=1.2)
        top += 1.4

    output = "Presentacion_Camiones_Vigilancia_iTTel.pptx"
    prs.save(output)
    print(f"✅ Presentación generada: {output}")
    return output


if __name__ == "__main__":
    try:
        generar_presentacion_camiones()
    except Exception as e:
        print("❌ Error al generar la presentación:", e)
        print("💡 Instala dependencias: pip install python-pptx Pillow")
