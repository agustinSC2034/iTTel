"""
Genera una presentación PowerPoint NUEVA y específica para la propuesta comercial
de Videovigilancia Móvil para la flota de recolección de residuos del GCBA.

Esta es una propuesta comercial independiente, NO modifica presentaciones existentes.
Incluye experiencia en VITTAL y contenido técnico del sistema MDVR.

Requiere: python-pptx, Pillow
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


# Paleta corporativa iTTel (desde assets/css/main.css)
COLOR_PRIMARY = RGBColor(34, 125, 179)   # #227db3
COLOR_ACCENT = RGBColor(56, 189, 248)    # #38bdf8
COLOR_TEXT_DARK = RGBColor(26, 32, 44)   # #1a202c
COLOR_TEXT_MEDIUM = RGBColor(74, 85, 104)  # #4a5568
COLOR_BG_LIGHT = RGBColor(248, 250, 252)  # #f8fafc
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_BORDER = RGBColor(229, 231, 235)


def add_logo(slide, logo_path):
    """Agrega el logo iTTel en la esquina superior izquierda."""
    if os.path.exists(logo_path):
        try:
            slide.shapes.add_picture(logo_path, Inches(0.4), Inches(0.4), width=Inches(1.8))
        except Exception:
            pass


def add_section_title(slide, title, subtitle=None):
    """Título de sección con subtítulo opcional."""
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.6), Inches(8.6), Inches(0.7))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_DARK
    p.alignment = PP_ALIGN.LEFT
    
    if subtitle:
        tb2 = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(8.6), Inches(0.5))
        tf2 = tb2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(16)
        p2.font.color.rgb = COLOR_TEXT_MEDIUM
        p2.alignment = PP_ALIGN.LEFT


def add_bullets(slide, items, left=0.8, top=2.0, width=8.4, height=4.8, font_size=17):
    """Agrega una lista con bullets."""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = COLOR_TEXT_DARK
        p.line_spacing = 1.4
        p.space_before = Pt(6) if i > 0 else Pt(0)


def add_card_box(slide, title, content, left, top, width, height):
    """Crea una tarjeta con título y contenido."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_WHITE
    shape.line.color.rgb = COLOR_BORDER
    shape.line.width = Pt(1.5)
    
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.margin_right = Inches(0.25)
    tf.margin_top = Inches(0.2)
    tf.margin_bottom = Inches(0.2)
    
    # Título
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(19)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.alignment = PP_ALIGN.LEFT
    
    # Contenido
    p2 = tf.add_paragraph()
    p2.text = content
    p2.font.size = Pt(14)
    p2.font.color.rgb = COLOR_TEXT_MEDIUM
    p2.space_before = Pt(8)
    p2.line_spacing = 1.3


def safe_add_picture(slide, img_path, left, top, width=None, height=None):
    """Agrega una imagen si existe."""
    if os.path.exists(img_path):
        try:
            if width:
                slide.shapes.add_picture(img_path, Inches(left), Inches(top), width=Inches(width))
            elif height:
                slide.shapes.add_picture(img_path, Inches(left), Inches(top), height=Inches(height))
            else:
                slide.shapes.add_picture(img_path, Inches(left), Inches(top))
        except Exception:
            pass


def generar_propuesta_videovigilancia():
    """Genera la presentación completa de la propuesta comercial."""
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Rutas de assets
    logo_path = os.path.join("assets", "images", "logoIttel1.jpg")
    camiones_dir = "presentacion_camiones"
    assets_images = os.path.join("assets", "images")
    
    # ============================================================
    # SLIDE 1: PORTADA
    # ============================================================
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s1.background.fill
    bg.solid()
    bg.fore_color.rgb = COLOR_PRIMARY
    
    # Logo grande en portada
    if os.path.exists(logo_path):
        try:
            s1.shapes.add_picture(logo_path, Inches(3.5), Inches(1.0), width=Inches(3.0))
        except Exception:
            pass
    
    # Título principal
    tb_title = s1.shapes.add_textbox(Inches(1.0), Inches(3.0), Inches(8.0), Inches(1.2))
    tf = tb_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Propuesta Comercial"
    p.font.size = Pt(50)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtítulo
    tb_sub = s1.shapes.add_textbox(Inches(1.0), Inches(4.2), Inches(8.0), Inches(1.5))
    tf2 = tb_sub.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "Sistema de Videovigilancia Móvil\npara Flota de Recolección de Residuos"
    p2.font.size = Pt(28)
    p2.font.color.rgb = COLOR_WHITE
    p2.alignment = PP_ALIGN.CENTER
    p2.line_spacing = 1.3
    
    # Cliente
    tb_client = s1.shapes.add_textbox(Inches(1.0), Inches(6.0), Inches(8.0), Inches(0.6))
    tf3 = tb_client.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "Gobierno de la Ciudad de Buenos Aires"
    p3.font.size = Pt(22)
    p3.font.color.rgb = RGBColor(200, 230, 255)
    p3.alignment = PP_ALIGN.CENTER
    
    # ============================================================
    # SLIDE 2: ÍNDICE / AGENDA
    # ============================================================
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_logo(s2, logo_path)
    add_section_title(s2, "Agenda de la Propuesta")
    
    agenda_items = [
        "1. Resumen Ejecutivo",
        "2. El Desafío: Necesidades del GCBA",
        "3. La Solución Propuesta: Sistema MDVR Completo",
        "4. Componentes Clave del Sistema",
        "5. Integración y Operación Automática",
        "6. Casos de Uso Operativos",
        "7. Arquitectura de Comunicación",
        "8. Plan de Implementación",
        "9. Beneficios para el GCBA",
        "10. Nuestra Experiencia: Proyecto VITTAL",
        "11. Próximos Pasos",
    ]
    add_bullets(s2, agenda_items, top=1.8, font_size=16)
    
    # ============================================================
    # SLIDE 3: RESUMEN EJECUTIVO
    # ============================================================
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_logo(s3, logo_path)
    add_section_title(s3, "Resumen Ejecutivo",
                      "Monitoreo, auditoría y seguridad 24/7 para la flota de residuos")
    
    bullets = [
        "✓ Objetivo: visibilidad total de la operación de recolección con auditoría instantánea.",
        "✓ Solución: MDVR vehicular + cámaras IP67/IK10 + 4G/GPS + software centralizado.",
        "✓ Auditoría acelerada: sensor en brazo hidráulico genera eventos automáticos.",
        "✓ Seguridad incrementada: G-sensor detecta colisiones y protege evidencia.",
        "✓ Equipamiento de grado industrial: resistente a vibración, agua y polvo.",
        "✓ Operación 100% automática y soporte 24×7×365.",
    ]
    add_bullets(s3, bullets, top=2.1, font_size=16)
    
    # Imagen ilustrativa
    safe_add_picture(s3, os.path.join(camiones_dir, "camion2.png"), 6.8, 2.2, height=3.5)
    
    # ============================================================
    # SLIDE 4: EL DESAFÍO
    # ============================================================
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_logo(s4, logo_path)
    add_section_title(s4, "El Desafío",
                      "¿Por qué necesita el GCBA este sistema?")
    
    desafios = [
        "🎯 Auditar el proceso de levantamiento de contenedores ante reclamos ciudadanos.",
        "⚖️ Generar evidencia irrefutable en incidentes viales y operativos.",
        "📍 Gestionar ubicación, eventos y estado de la flota en tiempo real.",
        "💪 Entorno hostil: vibración constante, agua, polvo y operación ruda 24/7.",
        "⏱️ Búsqueda manual en horas de grabación es inviable: se necesita búsqueda por eventos.",
    ]
    add_bullets(s4, desafios, top=2.0, font_size=17)
    
    # ============================================================
    # SLIDE 5: LA SOLUCIÓN PROPUESTA
    # ============================================================
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_logo(s5, logo_path)
    add_section_title(s5, "La Solución Propuesta",
                      "Ecosistema completo de videovigilancia móvil")
    
    add_card_box(s5, "1. Hardware Embarcado (MDVR)",
                 "Grabación multi-canal, módulos 4G/GPS, G-sensor integrado, supercapacitores y SSD industrial.",
                 0.6, 2.0, 2.8, 1.7)
    
    add_card_box(s5, "2. Red de Comunicación",
                 "Transmisión en vivo por 4G/LTE; telemetría GPS para geoposición en mapa.",
                 3.6, 2.0, 2.8, 1.7)
    
    add_card_box(s5, "3. Software de Gestión",
                 "Hik-Central Professional: mapa GIS, live view, búsqueda por eventos, protección automática de evidencia.",
                 6.6, 2.0, 2.8, 1.7)
    
    add_card_box(s5, "4. Sensores y Auditoría",
                 "Sensor en brazo hidráulico genera TAG en cada levantamiento para auditoría instantánea.",
                 0.6, 4.0, 2.8, 1.7)
    
    add_card_box(s5, "5. Seguridad Vial",
                 "G-sensor detecta colisiones/frenadas bruscas y protege automáticamente la evidencia.",
                 3.6, 4.0, 2.8, 1.7)
    
    add_card_box(s5, "6. Operación 24×7×365",
                 "NOC dedicado, monitoreo proactivo, soporte técnico y mantenimiento preventivo.",
                 6.6, 4.0, 2.8, 1.7)
    
    # ============================================================
    # SLIDE 6: EQUIPAMIENTO - MDVR
    # ============================================================
    s6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_logo(s6, logo_path)
    add_section_title(s6, "Equipamiento Principal: MDVR",
                      "Hikvision DS-MP5604H-S/GW — El cerebro del sistema")
    
    specs = [
        "✓ Grabador híbrido vehicular (analógico/IP).",
        "✓ Módulos integrados: 4G/LTE (chip SIM) y GPS.",
        "✓ G-Sensor interno (acelerómetro 3 ejes): detecta colisiones, frenadas y vuelcos.",
        "✓ Supercapacitores: cierre seguro de archivos ante pérdida súbita de batería.",
        "✓ Bahía 2.5\" con llave para SSD industrial (protección de evidencia).",
        "✓ Entradas de alarma (I/O) para sensores externos.",
    ]
    add_bullets(s6, specs, left=0.8, top=2.0, width=4.8, font_size=15)
    
    # Imagen del MDVR
    safe_add_picture(s6, os.path.join(camiones_dir, "MVR DS-MP5604-SD.png"), 5.8, 1.8, height=4.2)
    
    # ============================================================
    # SLIDE 7: ALMACENAMIENTO - SSD
    # ============================================================
    s7 = prs.slides.add_slide(prs.slide_layouts[6])
    add_logo(s7, logo_path)
    add_section_title(s7, "Almacenamiento: SSD Industrial",
                      "Clave para resistir el entorno vehicular")
    
    ssd_bullets = [
        "✓ SSD de grado industrial (no HDD ni tarjetas SD).",
        "✓ Resistencia extrema a vibración y golpes continuos.",
        "✓ Mayor confiabilidad de datos y durabilidad.",
        "✓ Protección contra corrupción de archivos.",
        "✓ Bahía con llave en MDVR para seguridad física de evidencia.",
    ]
    add_bullets(s7, ssd_bullets, top=2.0, font_size=17)
    
    # ============================================================
    # SLIDE 8: CÁMARAS Y CONECTORES
    # ============================================================
    s8 = prs.slides.add_slide(prs.slide_layouts[6])
    add_logo(s8, logo_path)
    add_section_title(s8, "Cámaras y Conectores",
                      "Diseñadas para entorno vehicular hostil")
    
    cam_bullets = [
        "✓ Cámaras Hikvision serie DS-2CS… (ej. DS-2CS54D7T-M).",
        "✓ Certificación IK10 (antivandálicas) e IP67 (herméticas: agua/polvo).",
        "✓ Conectores de aviación M12 roscados: garantizan conexión permanente ante vibración.",
        "✓ Múltiples canales: frontal, lateral, trasera, cabina, zona de carga.",
    ]
    add_bullets(s8, cam_bullets, left=0.8, top=2.0, width=4.6, font_size=16)
    
    # Imagen de cámara
    safe_add_picture(s8, os.path.join(camiones_dir, "camara_domo.png"), 5.8, 1.8, height=4.0)
    
    # ============================================================
    # SLIDE 9: INTEGRACIÓN ELÉCTRICA
    # ============================================================
    s9 = prs.slides.add_slide(prs.slide_layouts[6])
    add_logo(s9, logo_path)
    add_section_title(s9, "Integración Eléctrica e Ignición",
                      "Operación 100% automática, sin intervención del conductor")
    
    elec_bullets = [
        "✓ MDVR instalado oculto en cabina (bajo asiento o tablero).",
        "✓ Cableado B+ (batería directa), GND (tierra), IGN (contacto de llave).",
        "✓ Al girar la llave: MDVR arranca y comienza grabación automáticamente.",
        "✓ Al apagar: graba por tiempo programado (ej. 15 min) y entra en stand-by.",
        "✓ No agota batería del vehículo, consumo mínimo en reposo.",
    ]
    add_bullets(s9, elec_bullets, top=2.0, font_size=16)
    
    # ============================================================
    # SLIDE 10: SENSOR DE BRAZO Y AUDITORÍA
    # ============================================================
    s10 = prs.slides.add_slide(prs.slide_layouts[6])
    add_logo(s10, logo_path)
    add_section_title(s10, "Auditoría: Sensor en Brazo Hidráulico",
                       "De horas de búsqueda manual a segundos de auditoría")
    
    audit_bullets = [
        "✓ Sensor externo (magnético o inductivo) instalado en brazo hidráulico.",
        "✓ Cableado a Entrada de Alarma (I/O) del MDVR.",
        "✓ Cada activación del brazo genera una marca de evento (TAG) en la grabación.",
        "✓ Operador filtra en Hik-Central: 'mostrar todos los eventos de Entrada 1 entre X y Y horas'.",
        "✓ Resultado: lista de clips de 30 segundos, solo momentos de levantamiento de contenedores.",
        "✓ Auditoría que tomaría horas se resuelve en segundos.",
    ]
    add_bullets(s10, audit_bullets, top=2.0, font_size=15)
    
    # Imagen de monitoreo de alarmas
    safe_add_picture(s10, os.path.join(camiones_dir, "AlarmManagement.png"), 6.5, 2.0, height=3.5)
    
    # ============================================================
    # SLIDE 11: PLATAFORMA DE MONITOREO
    # ============================================================
    s11 = prs.slides.add_slide(prs.slide_layouts[6])
    add_logo(s11, logo_path)
    add_section_title(s11, "Plataforma de Monitoreo: Hik-Central",
                       "Control total de la flota desde el NOC")
    
    plat_bullets = [
        "✓ Mapa GIS en tiempo real con posición GPS de todos los camiones.",
        "✓ Live view por 4G/LTE: cámaras en vivo desde cualquier camión.",
        "✓ Reproducción centralizada y búsqueda por fecha/hora/evento.",
        "✓ Protección automática de evidencia ante eventos críticos (G-sensor).",
        "✓ Roles y permisos por perfil de usuario.",
        "✓ Clientes web y apps móviles (Android/iOS).",
    ]
    add_bullets(s11, plat_bullets, left=0.8, top=2.0, width=4.6, font_size=15)
    
    # Imagen de plataforma
    safe_add_picture(s11, os.path.join(camiones_dir, "HikCentral.png"), 5.8, 1.9, height=3.8)
    
    # ============================================================
    # SLIDE 12: CASOS DE USO
    # ============================================================
    s12 = prs.slides.add_slide(prs.slide_layouts[6])
    add_logo(s12, logo_path)
    add_section_title(s12, "Casos de Uso Operativos",
                       "¿Cómo se usa el sistema día a día?")
    
    add_card_box(s12, "Caso 1: Monitoreo en Vivo",
                 "Operador abre mapa GIS, visualiza ubicación de camiones y accede a cámaras en tiempo real por 4G.",
                 0.6, 2.0, 2.9, 1.7)
    
    add_card_box(s12, "Caso 2: Auditoría por Reclamo",
                 "Vecino reclama: 'camión rompió mi contenedor a las 10 AM'. Operador filtra por TAG de sensor de brazo → clips de 30 s → auditoría en segundos.",
                 3.6, 2.0, 2.9, 1.7)
    
    add_card_box(s12, "Caso 3: Incidente Vial",
                 "Camión tiene colisión → G-sensor detecta impacto → MDVR marca evento y protege video → evidencia forense asegurada.",
                 6.6, 2.0, 2.9, 1.7)
    
    add_card_box(s12, "Caso 4: Análisis de Rutas",
                 "GPS tracking permite analizar rutas, optimizar logística y validar cumplimiento de recorridos.",
                 0.6, 4.0, 2.9, 1.7)
    
    add_card_box(s12, "Caso 5: Capacitación",
                 "Grabaciones se usan para entrenar nuevos operadores y mejorar prácticas de conducción segura.",
                 3.6, 4.0, 2.9, 1.7)
    
    add_card_box(s12, "Caso 6: Reportes Periódicos",
                 "Sistema genera reportes automáticos de eventos, incidentes y KPIs operativos.",
                 6.6, 4.0, 2.9, 1.7)
    
    # ============================================================
    # SLIDE 13: ARQUITECTURA DE COMUNICACIÓN
    # ============================================================
    s13 = prs.slides.add_slide(prs.slide_layouts[6])
    add_logo(s13, logo_path)
    add_section_title(s13, "Arquitectura de Comunicación",
                       "Del camión al centro de datos")
    
    arch_bullets = [
        "1️⃣ Camión: Cámaras → MDVR → Módulo 4G/LTE + GPS.",
        "2️⃣ Red móvil: Transmisión de video en vivo y telemetría.",
        "3️⃣ Centro de Datos: Servidor Hik-Central con base de datos y storage.",
        "4️⃣ NOC / Operadores: Acceso desde PCs, tablets y smartphones.",
        "5️⃣ Operación 24×7×365: Monitoreo proactivo, alertas y reportes.",
    ]
    add_bullets(s13, arch_bullets, top=2.0, font_size=17)
    
    # Imagen de tracking en tiempo real
    safe_add_picture(s13, os.path.join(camiones_dir, "Real-time_tracking.png"), 6.2, 2.0, height=3.8)
    
    # ============================================================
    # SLIDE 14: PLAN DE IMPLEMENTACIÓN
    # ============================================================
    s14 = prs.slides.add_slide(prs.slide_layouts[6])
    add_logo(s14, logo_path)
    add_section_title(s14, "Plan de Implementación",
                       "Roadmap del proyecto")
    
    plan_bullets = [
        "Fase 1: Relevamiento técnico y diseño (1–2 semanas)",
        "  • Site survey en camiones tipo",
        "  • Definición de ubicación de cámaras y MDVR",
        "  • HLD y LLD de la solución",
        "",
        "Fase 2: Piloto (2–4 semanas)",
        "  • Instalación en X camiones seleccionados",
        "  • Pruebas funcionales y UAT",
        "  • Ajustes y optimización",
        "",
        "Fase 3: Despliegue masivo (escala progresiva)",
        "  • Instalación por lotes en toda la flota",
        "  • Capacitación de operadores y choferes",
        "",
        "Fase 4: Operación y Soporte",
        "  • Puesta en producción",
        "  • Soporte 24×7×365 desde NOC",
        "  • Mantenimiento preventivo y correctivo",
    ]
    add_bullets(s14, plan_bullets, top=2.0, font_size=14)
    
    # ============================================================
    # SLIDE 15: BENEFICIOS PARA GCBA
    # ============================================================
    s15 = prs.slides.add_slide(prs.slide_layouts[6])
    add_logo(s15, logo_path)
    add_section_title(s15, "Beneficios para el GCBA",
                       "¿Qué gana el Gobierno con esta inversión?")
    
    benefits = [
        "✅ Mejora de SLA y trazabilidad del servicio de recolección.",
        "✅ Reducción drástica de tiempos de auditoría (de horas a segundos).",
        "✅ Disminución de reclamos ciudadanos gracias a evidencia objetiva.",
        "✅ Seguridad operativa y vial incrementada (detección de incidentes).",
        "✅ Evidencia forense confiable y protegida ante litigios.",
        "✅ Datos para analítica, optimización de rutas y mejora continua.",
        "✅ Transparencia y rendición de cuentas ante la ciudadanía.",
    ]
    add_bullets(s15, benefits, top=2.0, font_size=17)
    
    # ============================================================
    # SLIDE 16: NUESTRA EXPERIENCIA - VITTAL
    # ============================================================
    s16 = prs.slides.add_slide(prs.slide_layouts[6])
    add_logo(s16, logo_path)
    add_section_title(s16, "Nuestra Experiencia: Proyecto VITTAL",
                       "Sistema de videovigilancia IP centralizado")
    
    vittal_bullets = [
        "✓ Cliente: VITTAL (empresa de emergencias médicas).",
        "✓ Solución: Plataforma de videovigilancia offline y online.",
        "✓ Características implementadas:",
        "  • Visualización y grabación centralizada de todas las cámaras",
        "  • Control remoto de cámaras (PTZ)",
        "  • Soporte de visión nocturna, detección por movimiento, 4K",
        "  • Alarmas y zonas de grabación configurables",
        "  • Apps Android/iOS y acceso web",
        "✓ Resultado: Sistema robusto, escalable y fácil de operar.",
        "✓ Experiencia aplicable: arquitectura similar para flota móvil.",
    ]
    add_bullets(s16, vittal_bullets, left=0.8, top=2.0, width=4.8, font_size=15)
    
    # Imagen del proyecto VITTAL
    safe_add_picture(s16, os.path.join(assets_images, "proyectos", "viTTal.png"), 5.8, 2.0, height=3.5)
    
    # ============================================================
    # SLIDE 17: SLA Y SOPORTE
    # ============================================================
    s17 = prs.slides.add_slide(prs.slide_layouts[6])
    add_logo(s17, logo_path)
    add_section_title(s17, "SLA y Soporte",
                       "Compromiso de servicio 24×7×365")
    
    sla_bullets = [
        "✓ Atención 24×7×365 desde NOC iTTel.",
        "✓ Monitoreo proactivo de la plataforma y conectividad.",
        "✓ KPIs operativos y reportes periódicos.",
        "✓ Mantenimiento preventivo programado.",
        "✓ Mantenimiento correctivo con SLA de respuesta.",
        "✓ Stock de partes críticas y reemplazo prioritario.",
        "✓ Actualizaciones de software y firmware.",
    ]
    add_bullets(s17, sla_bullets, top=2.0, font_size=17)
    
    # ============================================================
    # SLIDE 18: PRÓXIMOS PASOS
    # ============================================================
    s18 = prs.slides.add_slide(prs.slide_layouts[6])
    add_logo(s18, logo_path)
    add_section_title(s18, "Próximos Pasos",
                       "¿Cómo avanzamos?")
    
    next_bullets = [
        "1. Workshop técnico: validar alcance, requerimientos y expectativas.",
        "2. Definir cantidad de camiones, canales por vehículo y configuración.",
        "3. Site survey en camiones tipo para diseño detallado.",
        "4. Elaboración de HLD/LLD y cronograma de proyecto.",
        "5. Propuesta económica detallada.",
        "6. Aprobación y firma de contrato.",
        "7. Kick-off de proyecto y fase piloto.",
    ]
    add_bullets(s18, next_bullets, top=2.0, font_size=17)
    
    # ============================================================
    # SLIDE 19: CONTACTO
    # ============================================================
    s19 = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s19.background.fill
    bg.solid()
    bg.fore_color.rgb = COLOR_BG_LIGHT
    
    # Título
    tb = s19.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(8.0), Inches(0.8))
    p = tb.text_frame.paragraphs[0]
    p.text = "Contacto"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_DARK
    p.alignment = PP_ALIGN.CENTER
    
    # Subtítulo
    tb2 = s19.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(8.0), Inches(0.5))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = "Estamos a su disposición para ampliar cualquier aspecto de esta propuesta"
    p2.font.size = Pt(16)
    p2.font.color.rgb = COLOR_TEXT_MEDIUM
    p2.alignment = PP_ALIGN.CENTER
    
    # Tarjetas de contacto
    add_card_box(s19, "📞 Teléfonos",
                 "011-5272-9000\n0810-345-ITTEL (4883)",
                 0.9, 3.0, 2.7, 1.4)
    
    add_card_box(s19, "📧 Emails",
                 "administracion@it-tel.com.ar\nsoporte@it-tel.com.ar",
                 3.7, 3.0, 2.7, 1.4)
    
    add_card_box(s19, "💼 LinkedIn",
                 "/grupoittel",
                 6.5, 3.0, 2.7, 1.4)
    
    add_card_box(s19, "📍 Dirección",
                 "Av. Alicia Moreau de Justo 1930\nCiudad de Buenos Aires",
                 0.9, 4.6, 2.7, 1.4)
    
    add_card_box(s19, "🌐 Web",
                 "www.it-tel.com.ar",
                 3.7, 4.6, 2.7, 1.4)
    
    add_card_box(s19, "⏰ Soporte",
                 "24×7×365\nNOC Operativo",
                 6.5, 4.6, 2.7, 1.4)
    
    # Logo al pie
    if os.path.exists(logo_path):
        try:
            s19.shapes.add_picture(logo_path, Inches(4.0), Inches(6.4), width=Inches(2.0))
        except Exception:
            pass
    
    # Guardar presentación
    output = "Propuesta_Videovigilancia_Camiones_GCBA.pptx"
    prs.save(output)
    print(f"✅ Propuesta comercial generada: {output}")
    print(f"📊 Total de slides: {len(prs.slides)}")
    return output


if __name__ == "__main__":
    try:
        generar_propuesta_videovigilancia()
        print("\n✨ ¡Presentación lista para presentar al GCBA!")
    except Exception as e:
        print(f"❌ Error: {e}")
        print("💡 Verifica que estén instaladas las dependencias: pip install python-pptx Pillow")
